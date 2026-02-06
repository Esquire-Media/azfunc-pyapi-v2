from azure.durable_functions import Blueprint
from azure.storage.blob import BlobClient
from io import BytesIO
from libs.utils.azure_storage import get_blob_sas, get_container_client, init_blob_client
from libs.utils.pptx import add_custom_image, replace_text
from libs.utils.esquire.onspot_graphics.demographics import Demographics
from libs.utils.esquire.onspot_graphics.observations import Observations
from libs.utils.esquire.onspot_graphics.slider import SliderGraph
from pptx import Presentation
import os, pandas as pd, re
from libs.utils.pptx_helpers import duplicate_slide, get_shape, remove_shape

# Create a Blueprint instance for defining Azure Functions
bp = Blueprint()


# Define an activity function
@bp.activity_trigger(input_name="settings")
def activity_locationInsights_buildReport(settings: dict):

    # container for storing data and where the report will be exported
    runtime_container = get_container_client(
        connection_string=os.environ[settings["runtime_container"]["conn_str"]],
        container_name=settings["runtime_container"]["container_name"],
    )
    # container that holds static assets such as PPTX templates and image files
    resources_container = get_container_client(
        connection_string=os.environ[settings["resources_container"]["conn_str"]],
        container_name=settings["resources_container"]["container_name"],
    )

    template = Presentation(BytesIO(
        resources_container.download_blob(
            blob=f"templates/{settings['template']}.pptx"
        ).content_as_bytes()
    ))
    base_slide_count = len(template.slides)
    
    for location in settings["locations_data"]:
        # load location info from blob
        locationID = location["locationID"]

        # Load location info
        location_data = pd.read_csv(get_blob_sas(init_blob_client(
            conn_str=os.environ[settings["runtime_container"]["conn_str"]],
            container_name=settings["runtime_container"]["container_name"],
            blob_name=location["location_blob"],
        )))
        info = location_data.iloc[0].copy()
        info["Full Address"] = (
            f"{info['Address']}, {info['City']}, {info['State']} {info['Zip']}"
        )

        # load observations data from blob
        observations_data = pd.read_csv(
            get_blob_sas(
                init_blob_client(
                    conn_str=os.environ[settings["runtime_container"]["conn_str"]],
                    container_name=settings["runtime_container"]["container_name"],
                    blob_name=[
                        *runtime_container.list_blobs(
                            name_starts_with=f"{location['observations_blob']}/"
                        )
                    ][0]["name"],
                )
            )
        )
        observations = Observations(observations_data)


    # # load template presentation
    # template = Presentation(
    #     BytesIO(
    #         resources_container.download_blob(
    #             blob=f"templates/{settings['template']}.pptx"
    #         ).content_as_bytes()
    #     )
    # )

    # Duplicate slides for this location and track new slide indices
        # Track newly created slides and get their indices correctly
        new_slides = []
        for idx in range(base_slide_count):
            new_slide = duplicate_slide(template, idx)
            new_slides.append(template.slides[-1])  # Always get the last slide (newly added)

        # Ensure replacements happen only in duplicated slides
        replace_text_in_slides(slides=new_slides, location_info=info, observations=observations)

        execute_graphics_replacements(
            obs=observations, 
            demos=None, 
            template=template, 
            resources_client=resources_container, 
            settings=settings, 
            slides=new_slides
            )
        
        del observations
        # load demographics data from blob
        demographics_data = pd.read_csv(
            get_blob_sas(
                BlobClient.from_connection_string(
                    conn_str=os.environ[settings["runtime_container"]["conn_str"]],
                    container_name=settings["runtime_container"]["container_name"],
                    blob_name=[
                        *ContainerClient.from_connection_string(
                            conn_str=os.environ[settings["runtime_container"]["conn_str"]],
                            container_name=settings["runtime_container"]["container_name"],
                        ).list_blobs(
                            name_starts_with=f"{location['demographics_blob']}/"
                        )
                    ][0]["name"],
                )
            ),
            usecols = [
                "location (venue)",
                "state",
                "city",
                "zip",
                "zip4",
                "hashed device id",
                "number of instances (count seen at location)",
                "estimated_age 18-24",
                "estimated_age 25-29",
                "estimated_age 30-34",
                "estimated_age 35-39",
                "estimated_age 40-44",
                "estimated_age 45-49",
                "estimated_age 50-54",
                "estimated_age 55-59",
                "estimated_age 60-64",
                "estimated_age 65-69",
                "estimated_age 70+",
                "household_income under $15000",
                "household_income $15000 - $24999",
                "household_income $25000 - $34999",
                "household_income $35000 - $49999",
                "household_income $50000 - $74999",
                "household_income $75000 - $99999",
                "household_income $100000 - $149999",
                "household_income $150000 - $199999",
                "household_income $200000 - $249999",
                "household_income $250000+",
                "completed high school",
                "some college",
                "completed college",
                "attended vocational/technical school",
                "completed graduate school",
                "male",
                "female",
                "dwelling_type single family",
                "dwelling_type apt/multi-family",
                "married",
                "presence_of_children",
                "presence_of_veteran"
            ]
        )
        demographics = Demographics(demographics_data)

        execute_graphics_replacements(
            obs=None, 
            demos=demographics, 
            template=template, 
            resources_client=resources_container, 
            settings=settings, 
            slides=new_slides
            )
    
    # remove the blank template slides
    template = remove_slides_by_index(template, [0,1,2])

    # use a file-like object to export pptx as bytes
    output_io = BytesIO()
    template.save(output_io)
    output_io.seek(0)

    # upload bytes to blob storage
    filename = (
        re.sub(r"[^A-Za-z0-9 ]+", "", f"{settings['name']} {settings['batch_id']} {settings['report_id']}")
        .replace(" ", "_")
        .replace("__", "_")
    )
    output_blob_name = f"{settings['runtime_container']['output_blob']}/{filename}.pptx"
    runtime_container.upload_blob(name=output_blob_name, data=output_io, overwrite=True)

    return output_blob_name


def execute_graphics_replacements(
    obs: Observations,
    demos: Demographics,
    template: Presentation,
    resources_client: ContainerClient,
    settings: dict,
    slides
):
    
    """
    Create and export the device observations graphics.
    """
    traffic_slide = slides[1]
    demos_slide = slides[2]
    # heatmap_slide = template.slides[3]
    # creative_slide = template.slides[4]

    if obs is not None:
        # __TRAFFIC SLIDE__
        # foot traffic line graph
        add_custom_image(
            file=obs.foot_traffic_graph(return_bytes=True),
            slide=traffic_slide,
            placeholder=get_shape(traffic_slide, "Market Traffic"),
        )
        # time distribution graph
        add_custom_image(
            file=obs.time_distribution_graph(return_bytes=True),
            slide=traffic_slide,
            placeholder=get_shape(traffic_slide, "Traffic Distribution"),
        )
        # trend score slider
        add_custom_image(
            file=SliderGraph(
                val=obs.trend_score(),
                labels=["", "Decreasing", "", "Neutral", "", "Increasing", ""],
            ).export(return_bytes=True),
            slide=traffic_slide,
            placeholder=get_shape(traffic_slide, "Overall Trend_SliderPlaceholder"),
        )
        # stability score slider
        add_custom_image(
            file=SliderGraph(
                obs.stability_score(),
                labels=["", "Volatile", "", "Moderate", "", "Consistent", ""],
            ).export(return_bytes=True),
            slide=traffic_slide,
            placeholder=get_shape(traffic_slide, "Consistency_SliderPlaceholder"),
        )
        # recent performance slider
        add_custom_image(
            file=SliderGraph(
                val=obs.recent_score(),
                labels=["", "Below Average", "", "Average", "", "Above Average", ""],
            ).export(return_bytes=True),
            slide=traffic_slide,
            placeholder=get_shape(traffic_slide, "Recent Activity_SliderPlaceholder"),
        )
        # remove the placeholders
        remove_slider_placeholders(traffic_slide)

    if demos is not None:
        # --- DEMOGRAPHICS SLIDE ---
        add_custom_image(
            file=demos.gender_graph(return_bytes=True),
            slide=demos_slide,
            placeholder=get_shape(demos_slide, "GenderChart"),
        )

        add_custom_image(
            file=demos.marriage_graph(return_bytes=True),
            slide=demos_slide,
            placeholder=get_shape(demos_slide, "MarriageChart"),
        )

        add_custom_image(
            file=demos.dwelling_graph(return_bytes=True),
            slide=demos_slide,
            placeholder=get_shape(demos_slide, "DwellingChart"),
        )

        add_custom_image(
            file=demos.children_graph(return_bytes=True),
            slide=demos_slide,
            placeholder=get_shape(demos_slide, "ChildrenChart"),
        )

        add_custom_image(
            file=demos.age_graph(return_bytes=True),
            slide=demos_slide,
            placeholder=get_shape(demos_slide, "AgeChart"),
        )

        add_custom_image(
            file=demos.income_graph(return_bytes=True),
            slide=demos_slide,
            placeholder=get_shape(demos_slide, "IncomeChart"),
        )

    # # --HEATMAP SLIDE--
    # # heatmap scatter plot
    # add_custom_image(
    #     file=obs.heatmap_graph(return_bytes=True),
    #     slide=heatmap_slide,
    #     placeholder=heatmap_slide.shapes[11],
    # )
    # # promotional ad image
    # blob_client = resources_client.get_blob_client(
    #     blob=f"promotions/{settings['promotionSet']}/Ad.png"
    # )
    # add_custom_image(
    #     file=BytesIO(blob_client.download_blob().content_as_bytes()),
    #     slide=heatmap_slide,
    #     placeholder=heatmap_slide.shapes[12],
    # )
    # # promotional ad hyperlink
    # blob_client = resources_client.get_blob_client(
    #     blob=f"promotions/{settings['promotionSet']}/Ad Link.txt"
    # )
    # ad = heatmap_slide.shapes[-1]
    # ad.click_action.hyperlink.address = blob_client.download_blob().content_as_text()

    # # --CREATIVE SLIDE--
    # # placeholder indexes for each of the 3 creative images
    # creative_placeholders = {0: 18, 1: 19, 2: 20}
    # for i, blob_name in enumerate(
    #     [
    #         blob.name
    #         for blob in resources_client.list_blobs(
    #             f"creatives/{settings['creativeSet']}/"
    #         )
    #     ][:3]
    # ):
    #     # creative images
    #     blob_client = resources_client.get_blob_client(blob=blob_name)
    #     if blob_client.exists():
    #         add_custom_image(
    #             file=BytesIO(blob_client.download_blob().content_as_bytes()),
    #             slide=creative_slide,
    #             placeholder=creative_slide.shapes[creative_placeholders[i]],
    #         )


def replace_text_in_slides(slides, location_info, observations):
    """Replace placeholders with real data only in the duplicated slides."""
    # bullet points and text data
    latest_week = observations.get_latest_week()
    best_week = observations.get_best_week()
    worst_week = observations.get_worst_week()
    
    replacements = {
        "{{title}}": location_info["Owner"].upper(),
        "{{subtitle}}": location_info["Full Address"],
        "{{year}}": latest_week["Year"],
        "{{latest week}}": latest_week["Week"],
        "{{latest performance}}": latest_week["Performance"],
        "{{latest range}}": latest_week["Range"],
        "{{best week}}": best_week["Week"],
        "{{best performance}}": best_week["Performance"],
        "{{best range}}": best_week["Range"],
        "{{worst week}}": worst_week["Week"],
        "{{worst performance}}": worst_week["Performance"],
        "{{worst range}}": worst_week["Range"],
        "{{bullet1}}": observations.bullet_current_performance(),
        "{{bullet2}}": observations.bullet_continuous_growth(),
        "{{bullet3}}": observations.bullet_six_weeks(),
        "{{bullet4}}": observations.bullet_budget(),
    }
    from libs.utils.pptx import add_custom_image, replace_text
    # replace text in slides with generated stats and bullet points
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    replace_text(replacements, shapes)

def remove_slides_by_index(prs, removal_indices):
    # Sort indices in reverse order to avoid shifting issues
    for i in sorted(removal_indices, reverse=True):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    return prs

# remove slider placeholders (important to do this last and in decreasing order to not mess up other shape indexes)
def remove_slider_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.name.endswith("SliderPlaceholder"):
            remove_shape(shape)
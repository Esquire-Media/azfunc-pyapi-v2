from azure.durable_functions import Blueprint
from azure.storage.blob import ContainerClient
from datetime import date
from io import BytesIO
from libs.utils.text import pad_text
from libs.utils.pptx import replace_text, add_custom_image
from libs.utils.azure_storage import get_blob_sas
from pptx import Presentation
import pandas as pd, os

# Create a Blueprint instance for defining Azure Functions
bp = Blueprint()

# Define an activity function
@bp.activity_trigger(input_name="settings")
def activity_campaignProposal_executeReport(settings: dict):

    # TEMPLATE IMPORT
    resources_client: ContainerClient = ContainerClient.from_connection_string(conn_str=os.environ[settings["resources_container"]['conn_str']], container_name=settings["resources_container"]["container_name"])
    blob_client = resources_client.get_blob_client(blob=f"templates/template.pptx")
    template = Presentation(BytesIO(blob_client.download_blob().content_as_bytes()))
    
    # DATA IMPORTS
    # import addresses
    container_client: ContainerClient = ContainerClient.from_connection_string(conn_str=os.environ[settings["runtime_container"]['conn_str']], container_name=settings["runtime_container"]["container_name"])
    blob_client = container_client.get_blob_client(blob=f"{settings['instance_id']}/addresses.csv")
    addresses = pd.read_csv(get_blob_sas(blob_client))

    # if we have optionalSlides being put in, then make sure we're not doing extra work
    mover_counts = pd.DataFrame()
    mover_totals = pd.DataFrame()
    if ('new_mover' in settings.get('optionalSlides', [])) or ('optionalSlides' not in settings.keys()):
        # import mover totals (NOTE: as a single Pandas row, not a Dataframe)
        blob_client = container_client.get_blob_client(blob=f"{settings['instance_id']}/mover_totals.csv")
        mover_totals = pd.read_csv(get_blob_sas(blob_client)).iloc[0]
        
        # import mover counts
        blob_client = container_client.get_blob_client(blob=f"{settings['instance_id']}/mover_counts.csv")
        mover_counts = pd.read_csv(get_blob_sas(blob_client))

    competitors = pd.DataFrame()
    if ('in_market_shopper' in settings.get('optionalSlides', [])) or ('optionalSlides' not in settings.keys()):
        # import competitor list
        blob_client = container_client.get_blob_client(blob=f"{settings['instance_id']}/competitors.csv")
        competitors = pd.read_csv(get_blob_sas(blob_client))


    # populate presentation with generated text and graphics
    execute_text_replacements(template=template, settings=settings, mover_counts=mover_counts, mover_totals=mover_totals, addresses=addresses, competitors=competitors)
    execute_graphics_replacements(template=template, settings=settings, container_client=container_client, resources_client=resources_client)

    # remove the extraneous slides
    template = remove_excess_optionalSlides(template, settings)

    # FORMATTED PPTX UPLOAD
    # use a file-like object to export pptx as bytes
    output_io = BytesIO()
    template.save(output_io)
    output_io.seek(0)

    # upload bytes to blob storage
    blob_client = container_client.get_blob_client(blob=f"{settings['instance_id']}/CampaignProposal-{settings['name']}.pptx")
    blob_client.upload_blob(data=output_io)

    return {}

def execute_text_replacements(template:Presentation, settings:dict, mover_counts:pd.DataFrame, mover_totals:pd.DataFrame, addresses:pd.DataFrame, competitors:pd.DataFrame):
    
    # initialize the replacement set with placeholders for the optional ones
    replacements = {
        "{{request name}}"    : settings['name'],
        "{{month}}"           : date.today().strftime('%B'),
        "{{year}}"            : date.today().year,
        "{{count locations}}" : len(addresses),
        "{{owned list}}"      : '\n'.join(addresses['address'].unique()),
        "{{mover headers}}"   : "",
        "{{mover counts}}"    : "",
        "{{mover totals}}"    : "",
        "{{r1}}"              : "",
        "{{r2}}"              : "",
        "{{r3}}"              : "",
        "{{competitor list}}" : ""
    }

    if ('new_mover' in settings.get('optionalSlides', [])) or ('optionalSlides' not in settings.keys()):
        # TEXT FORMATTIING
        # mover headers, data, and a totals row
        radii = settings['moverRadii']
        max_num_length = len(str("{:,}".format(mover_totals[f'movers_{radii[2]}mi'])))
        mover_headers = f"{pad_text('Address',26)}\t{radii[0]} Mile\t{radii[1]} Mile\t{radii[2]} Mile"
        mover_lines = []
        for i, row in mover_counts.iterrows():
            # by-location counts
            mover_lines.append(
                f"""{pad_text(row['address'],26)}\t{pad_text(row[f'movers_{radii[0]}mi'],max_num_length)}\t{pad_text(row[f'movers_{radii[1]}mi'],max_num_length)}\t{pad_text(row[f'movers_{radii[2]}mi'],max_num_length)}"""
            )
        # unique mover counts for entire audience
        mover_totals_line = f"{pad_text('Total Adjusted for Overlap',26)}\t{pad_text(mover_totals[f'movers_{radii[0]}mi'],max_num_length)}\t{pad_text(mover_totals[f'movers_{radii[1]}mi'],max_num_length)}\t{pad_text(mover_totals[f'movers_{radii[2]}mi'],max_num_length)}"

        # TEXT REPLACEMENTS
        replacements["{{mover headers}}"]   = mover_headers
        replacements["{{mover counts}}"]    = "\n".join(mover_lines)
        replacements["{{mover totals}}"]    = mover_totals_line
        replacements["{{r1}}"]              = radii[0]
        replacements["{{r2}}"]              = radii[1]
        replacements["{{r3}}"]              = radii[2]
    
    if ('in_market_shopper' in settings.get('optionalSlides', [])) or ('optionalSlides' not in settings.keys()):
        replacements["{{competitor list}}"] = '\n'.join(competitors['chain_name'].value_counts().index.tolist()[:35])
    
    # replace text in slides with generated stats and bullet points
    slides = [slide for slide in template.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    replace_text(replacements, shapes)

def execute_graphics_replacements(template:Presentation, settings:dict, container_client:ContainerClient, resources_client:ContainerClient):
    
    # GRAPHICS REPLACEMENTS
    movers_slide = template.slides[9]
    competitors_slide = template.slides[10]
    display_social_slide = template.slides[4]
    ott_slide = template.slides[5]

    if ('new_mover' in settings.get('optionalSlides', [])) or ('optionalSlides' not in settings.keys()):
        radii = settings['moverRadii']

        # mover maps
        blob_client = container_client.get_blob_client(blob=f"{settings['instance_id']}/mover_map_{radii[0]}mi.png")
        add_custom_image(BytesIO(blob_client.download_blob().content_as_bytes()), movers_slide, movers_slide.shapes[30])

        blob_client = container_client.get_blob_client(blob=f"{settings['instance_id']}/mover_map_{radii[1]}mi.png")
        add_custom_image(BytesIO(blob_client.download_blob().content_as_bytes()), movers_slide, movers_slide.shapes[31])

        blob_client = container_client.get_blob_client(blob=f"{settings['instance_id']}/mover_map_{radii[2]}mi.png")
        add_custom_image(BytesIO(blob_client.download_blob().content_as_bytes()), movers_slide, movers_slide.shapes[29])

    if ('in_market_shopper' in settings.get('optionalSlides', [])) or ('optionalSlides' not in settings.keys()):
        # competitors map
        blob_client = container_client.get_blob_client(blob=f"{settings['instance_id']}/competitors.png")
        add_custom_image(BytesIO(blob_client.download_blob().content_as_bytes()), competitors_slide, competitors_slide.shapes[19])

    # CREATIVE SET REPLACEMENTS
    creativeSet = settings['creativeSet']

    # search for images with the creativeSet
    display_creatives = [*resources_client.list_blobs(f"creatives/{creativeSet}/Display/300x250/")][:1] + [*resources_client.list_blobs(f"creatives/{creativeSet}/Display/300x600/")][:1] + [*resources_client.list_blobs(f"creatives/{creativeSet}/Display/728x90/")][:1]
    social_creatives = [*resources_client.list_blobs(f"creatives/{creativeSet}/Social/Feed/")][:3] + [*resources_client.list_blobs(f"creatives/{creativeSet}/Social/Reel/")][:2]
    ott_creatives = [*resources_client.list_blobs(f"creatives/{creativeSet}/OTT/")][:2]

    # borrow from the default creativeSet where necessary
    if len(display_creatives) < 3:
        display_creatives = [*resources_client.list_blobs(f"creatives/Default/Display/300x250/")][:1] + [*resources_client.list_blobs(f"creatives/Default/Display/300x600/")][:1] + [*resources_client.list_blobs(f"creatives/Default/Display/728x90/")][:1]
    if len(social_creatives) < 5:
        social_creatives = [*resources_client.list_blobs(f"creatives/Default/Social/Feed/")][:3] + [*resources_client.list_blobs(f"creatives/Default/Social/Reel/")][:2]
    if len(ott_creatives) < 2:
        ott_creatives = [*resources_client.list_blobs(f"creatives/Default/OTT/")][:2]

    # display_creatives
    display_300x250_creatives = [dc for dc in display_creatives if '300x250' in dc.name]
    display_300x600_creatives = [dc for dc in display_creatives if '300x600' in dc.name]
    display_728x90_creatives = [dc for dc in display_creatives if '728x90' in dc.name]

    social_feed_creatives = [sc for sc in social_creatives if 'Feed' in sc.name]
    social_reel_creatives = [sc for sc in social_creatives if 'Reel' in sc.name]

    # display 300x250
    blob_client = resources_client.get_blob_client(blob=display_300x250_creatives[0].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=display_social_slide,
        placeholder=display_social_slide.shapes[16],
    )

    # display 300x600
    blob_client = resources_client.get_blob_client(blob=display_300x600_creatives[0].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=display_social_slide,
        placeholder=display_social_slide.shapes[17],
    )

    # display 728x90
    blob_client = resources_client.get_blob_client(blob=display_728x90_creatives[0].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=display_social_slide,
        placeholder=display_social_slide.shapes[15],
    )

    # social feed 1
    blob_client = resources_client.get_blob_client(blob=social_feed_creatives[0].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=display_social_slide,
        placeholder=display_social_slide.shapes[18],
    )

    # social feed 2
    blob_client = resources_client.get_blob_client(blob=social_feed_creatives[1].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=display_social_slide,
        placeholder=display_social_slide.shapes[19],
    )

    # social feed 3  
    blob_client = resources_client.get_blob_client(blob=social_feed_creatives[2].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=display_social_slide,
        placeholder=display_social_slide.shapes[20],
    )  
    
    # social reel 1
    blob_client = resources_client.get_blob_client(blob=social_reel_creatives[0].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=display_social_slide,
        placeholder=display_social_slide.shapes[22],
    )

    # social reel 2
    blob_client = resources_client.get_blob_client(blob=social_reel_creatives[1].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=display_social_slide,
        placeholder=display_social_slide.shapes[21],
    )

    # ott banner 1
    blob_client = resources_client.get_blob_client(blob=ott_creatives[0].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=ott_slide,
        placeholder=ott_slide.shapes[3],
    )

    # ott banner 2
    blob_client = resources_client.get_blob_client(blob=ott_creatives[1].name)
    add_custom_image(
        file=BytesIO(blob_client.download_blob().content_as_bytes()),
        slide=ott_slide,
        placeholder=ott_slide.shapes[4],
    )


def remove_excess_optionalSlides(prs: Presentation, settings:dict):
    """
    Removes slides from a PowerPoint presentation based on settings

    Parameters:
    prs (Presentation): The PowerPoint presentation object.
    settings: dict of the input json body. 
        should contains a key of 'optionalSlides'
        E.G.: {
            'optionalSlides':[
                'next_steps',
                'new_mover'
            ]
        }

    Returns:
    Presentation: The modified PowerPoint presentation.
    """

    # the indices of the slides that are optional (0-indexed)
    # any of these keys that are in the settings optionalSlides are retained
    optional_slide_ids = {
        'new_mover': 9,
        'in_market_shopper':10,
        'pricing':11,
        'next_steps':12
    }

    # null handling
    keep_slides = set(settings.get('optionalSlides', []))

    # Sort in reverse to avoid shifting issues
    removal_indices = sorted(
        [idx for key, idx in optional_slide_ids.items() if key not in keep_slides],
        reverse=True
    )

    prs = remove_slides_by_index(prs, removal_indices)

    return prs

def remove_slides_by_index(prs, removal_indices):
    # remove the slides via xml relationships
    for i in removal_indices: 
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    return prs